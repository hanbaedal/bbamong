# -*- coding: utf-8 -*-
"""빠몽이 예측 게임 — API-SPORTS 연동 후 시스템 구조/흐름/문제점 PPT"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "PPAMONG_예측게임_시스템흐름.pptx"

C_TITLE = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT = RGBColor(0xE1, 0x19, 0x36)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_OK = RGBColor(0x16, 0x65, 0x34)
C_WARN = RGBColor(0xB4, 0x53, 0x09)
C_BAD = RGBColor(0x9F, 0x12, 0x39)
C_ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)
C_BOX = RGBColor(0xEE, 0xF2, 0xFF)
FONT = "Malgun Gothic"


def set_run(run, text, size=18, bold=False, color=C_TITLE):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, text, 10, color=C_MUTED)


def slide_title(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12), Inches(0.6))
    r = tb.text_frame.paragraphs[0].add_run()
    set_run(r, title, 26, True, C_WHITE)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(1.12), Inches(12), Inches(0.4))
        sr = sb.text_frame.paragraphs[0].add_run()
        set_run(sr, subtitle, 13, color=C_MUTED)
    return slide


def add_bullets(slide, items, left=0.55, top=1.55, width=12.2, size=14, gap=0.36):
    y = top
    for item in items:
        box = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(gap))
        box.text_frame.word_wrap = True
        r = box.text_frame.paragraphs[0].add_run()
        set_run(r, item, size)
        y += gap


def add_table(slide, headers, rows, left=0.45, top=1.55, col_widths=None, font_size=10):
    cols = len(headers)
    rows_n = len(rows) + 1
    if col_widths is None:
        col_widths = [12.3 / cols] * cols

    table = slide.shapes.add_table(
        rows_n,
        cols,
        Inches(left),
        Inches(top),
        Inches(sum(col_widths)),
        Inches(0.38 + 0.38 * rows_n),
    ).table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TITLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = C_WHITE

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_ROW_ALT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = C_TITLE


def add_flow_boxes(slide, labels, top=1.7, box_h=0.7):
    n = len(labels)
    gap = 0.18
    total_w = 12.2
    box_w = (total_w - gap * (n - 1)) / n
    x = 0.5
    for i, label in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top),
            Inches(box_w),
            Inches(box_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_BOX
        shape.line.color.rgb = C_TITLE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        set_run(r, label, 11, True)
        if i < n - 1:
            arrow = slide.shapes.add_textbox(
                Inches(x + box_w - 0.05),
                Inches(top + 0.18),
                Inches(0.25),
                Inches(0.35),
            )
            ar = arrow.text_frame.paragraphs[0].add_run()
            set_run(ar, "→", 16, True, C_ACCENT)
        x += box_w + gap


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    s = slide_title(
        prs,
        "빠몽이 예측 게임 시스템 구조·흐름",
        "API-SPORTS 연동 이후 · 관리자 / 운영자 / 사용자 · 문제점 점검",
    )
    add_bullets(
        s,
        [
            "• 데이터: API-SPORTS Baseball (스코어보드·일정·상태)",
            "• 정산: 포인트 선택 × 고정 배당 (아웃1.2 / 1루1.5 / 2루3 / 3루10 / 홈런5)",
            "• 타석 제어: 운영자 수동 (예측 시작·중지·결과) + 비상 수동 모드",
            "• 작성 기준: 2026-07 코드 반영본",
        ],
        top=1.8,
        size=16,
        gap=0.45,
    )
    add_footer(s, "ppamong.com · First-Visit Replit")

    # 2. 전체 구조
    s = slide_title(prs, "1. 전체 시스템 구조", "외부 API → 중앙 백엔드 → 3개 클라이언트")
    add_flow_boxes(
        s,
        [
            "API-SPORTS\nKBO 스코어",
            "백엔드\n폴링 2.5초",
            "MongoDB\nMatch/Prediction",
            "WebSocket\n/ws/match",
            "관리자·운영자\n·사용자",
        ],
        top=1.7,
    )
    add_table(
        s,
        ["계층", "역할", "자동/수동"],
        [
            ["API-SPORTS", "오늘 경기·이닝·점수·종료 상태", "자동 폴링"],
            ["백엔드", "매핑·정산·광고 트리거·헬스", "자동 + 수동 API"],
            ["관리자", "경기 등록·API 연결·모니터링·비상모드", "수동 중심"],
            ["운영자", "예측 시작/중지/결과·광고 강제", "수동 (타석)"],
            ["사용자", "배팅·광고 시청·결과 수신", "앱/웹 UI"],
        ],
        top=2.8,
        col_widths=[2.2, 6.5, 3.5],
    )
    add_footer(s, "문자중계·타석 결과·볼카운트는 API-SPORTS에 없음")

    # 3. 역할 비교
    s = slide_title(prs, "2. 역할별 자동/수동 구분", "무엇이 바뀌고 무엇이 남았는가")
    add_table(
        s,
        ["기능", "관리자", "운영자", "사용자"],
        [
            ["오늘 경기 등록", "수동 필수", "-", "-"],
            ["API 경기 연결", "수동 필수(저장 후)", "-", "-"],
            ["이닝·점수 표시", "자동", "자동", "자동"],
            ["예측 시작/중지", "비상 시 가능", "수동 필수", "신호 수신"],
            ["결과(1루~아웃)", "비상 시 가능", "수동 필수", "정산 수신"],
            ["경기 종료", "강제 종료 유지", "-", "API 종료 시 자동"],
            ["공수교대 광고", "모니터링", "강제 시작/중지", "시청·500P"],
            ["배팅 금액", "-", "-", "포인트 선택"],
            ["정산", "-", "결과 입력 트리거", "amount×배당"],
        ],
        top=1.5,
        col_widths=[3.0, 3.1, 3.1, 3.1],
        font_size=11,
    )
    add_footer(s, "핵심: 스코어는 자동, 타석 판정은 여전히 사람")

    # 4. 관리자 흐름
    s = slide_title(prs, "3. 관리자 게임 흐름", "경기 관리 → API 연결 → 모니터링")
    add_flow_boxes(
        s,
        [
            "구장/경기\n등록",
            "저장하기",
            "API-SPORTS\n연결",
            "헬스 Green\n확인",
            "배팅·스코어\n모니터링",
            "비상 수동\n전환(필요시)",
        ],
        top=1.65,
        box_h=0.85,
    )
    add_bullets(
        s,
        [
            "① 관리자 → 경기 관리 → + 경기 등록 → 1~5경기 입력",
            "② 하단 「1. 저장하기」 → 「2. API-SPORTS 오늘 경기 연결」",
            "③ 실시간 게임 모니터링: 헬스 신호등, 스코어보드, 배팅 분포 그래프",
            "④ API 불안정 시 「수동 제어 전환」 → 자동 광고/자동 종료 중단, 운영자 수동 운영",
            "⑤ 광고 매출은 AdMob 리포트(기존) — API-SPORTS와 별개",
        ],
        top=2.8,
        size=14,
        gap=0.42,
    )
    add_footer(s, "연결 생략 시: 예측은 가능하나 스코어 자동·자동 종료 없음")

    # 5. 운영자 흐름
    s = slide_title(prs, "4. 운영자 게임 흐름", "타석 단위 제어는 운영자 책임")
    add_flow_boxes(
        s,
        [
            "배정 경기\n입장",
            "스코어보드\n확인",
            "예측 시작",
            "예측 중지",
            "결과 입력\n1루~아웃",
            "다음 라운드\n/광고",
        ],
        top=1.65,
        box_h=0.85,
    )
    add_bullets(
        s,
        [
            "• API가 이닝·점수를 보여주므로 ‘눈대중 스코어’ 부담은 감소",
            "• 타석 시작·종료·결과(아웃/1루/2루/3루/홈런)는 API에 없어 버튼 유지",
            "• 공수교대/투수교체: 라운드 강제 진행 + 광고 시작/중지 가능",
            "• 비상 수동 모드일 때: 중계를 보며 예측 종료·결과 정산을 강제로 이어감",
            "• 경기 최종 종료는 API 상태(FT 등)로 자동 처리(자동 모드 시)",
        ],
        top=2.8,
        size=14,
        gap=0.42,
    )
    add_footer(s, "운영자 = 타석 판정관 / API = 스코어보드 보조")

    # 6. 사용자 흐름
    s = slide_title(prs, "5. 사용자 게임 흐름", "선택 → 배팅 → 대기 → 정산 → 광고")
    add_flow_boxes(
        s,
        [
            "경기 선택",
            "금액 선택\n50~1000P",
            "예측 선택\n아웃~홈런",
            "결과 대기",
            "적중 시\n금액×배당",
            "공수교대\n광고·배너",
        ],
        top=1.65,
        box_h=0.85,
    )
    add_table(
        s,
        ["선택지", "배당", "예: 100P 적중 시"],
        [
            ["아웃", "1.2배", "120P"],
            ["1루", "1.5배", "150P"],
            ["2루", "3배", "300P"],
            ["3루", "10배", "1000P"],
            ["홈런", "5배", "500P"],
        ],
        top=2.85,
        col_widths=[3.0, 3.0, 4.0],
        left=1.5,
    )
    add_footer(s, "미적중 포인트는 소멸 · 기존 패자풀 분배 폐지")

    # 7. 라운드 E2E
    s = slide_title(prs, "6. 한 타석(라운드) 엔드투엔드", "백엔드 기준 시퀀스")
    add_bullets(
        s,
        [
            "1) 운영자 prediction/start → WS prediction_started → 유저 배팅 가능",
            "2) 유저 POST /predictions (amount 선택) → 즉시 차감",
            "3) 운영자 prediction/stop → WS prediction_stopped → 배팅 마감",
            "4) 운영자 result(아웃~홈런) → 고정배당 정산 → WS round_result",
            "5) 자동 nextRound → WS round_next + banner_ad_show (타자 교체 배너)",
            "6) API 이닝 변경 감지(자동모드) → ad_started + rewarded_ad_offer",
            "7) 유저 광고 완료 → /ad-reward 500P (5초 이내 취소 시 무보상)",
            "8) API 경기 종료(FT) → endMatch → WS match_ended",
        ],
        top=1.55,
        size=14,
        gap=0.4,
    )
    add_footer(s, "WebSocket path: /ws/match")

    # 8. 문제점
    s = slide_title(prs, "7. 문제점·리스크 점검", "운영 전 반드시 인지할 항목")
    add_table(
        s,
        ["등급", "문제", "영향", "대응"],
        [
            [
                "높음",
                "타석 결과 API 없음",
                "완전 무인 운영 불가",
                "운영자 수동 + 비상모드 유지",
            ],
            [
                "높음",
                "연결 누락",
                "스코어/자동종료 미동작",
                "등록 모달 2단계 버튼 강제 안내",
            ],
            [
                "높음",
                "매핑 순서 오류",
                "1~5경기 ↔ KBO 경기 뒤섞임",
                "시작시각 순 매핑 후 관리자 확인 UI 강화 필요",
            ],
            [
                "중간",
                "이닝변경=광고 휴리스틱",
                "오탐 시 광고 오발송",
                "운영자 광고 중지 + 수동모드",
            ],
            [
                "중간",
                "자동 종료 ≠ 미정산 라운드",
                "pending 예측 잔존 가능",
                "종료 전 라운드 마감·결과 입력 규약",
            ],
            [
                "중간",
                "API 지연(수 초~)",
                "중계보다 늦은 판정 체감",
                "운영자가 눈으로 최종 확정",
            ],
            [
                "낮음",
                "리그 ID/키 오류",
                "헬스 Red, 연결 실패",
                "Secrets·KBO league id 점검",
            ],
        ],
        top=1.45,
        col_widths=[1.2, 3.2, 3.4, 4.4],
        font_size=10,
    )
    add_footer(s, "빨강=운영 리스크 / 노랑=운영 규약으로 완화 / 초록=설정 이슈")

    # 9. 권장 운영 체크리스트
    s = slide_title(prs, "8. 권장 운영 체크리스트", "경기일 당일")
    add_bullets(
        s,
        [
            "□ Replit Secrets: API_SPORTS_KEY / API_SPORTS_KBO_LEAGUE_ID",
            "□ 관리자: 오늘 1~5경기 등록 → 저장 → API-SPORTS 연결",
            "□ 모니터링: 헬스 Green, 팀명·스코어 매핑 올바른지 확인",
            "□ 운영자 op1~op5: 배정 경기 입장, 스코어보드 표시 확인",
            "□ 예측 루프: 시작 → 유저 배팅 → 중지 → 결과 → (광고)",
            "□ API 이상 시: 관리자 수동 제어 전환 → 운영자 강제 종료/정산",
            "□ 경기 종료 직전: 열린 라운드 없는지 확인 후 API 자동종료 대기",
            "□ 배포: GitHub main push 후 Replit git pull + Deploy",
        ],
        top=1.55,
        size=15,
        gap=0.42,
    )
    add_footer(s, "다음 개선 후보: 연결 시 팀명 수동 매핑 UI, 미정산 라운드 종료 가드")

    # 10. 한 줄 요약
    s = slide_title(prs, "9. 한 줄 요약", "")
    add_bullets(
        s,
        [
            "• 관리자: 등록 + API 연결 + 감시(+비상)",
            "• 운영자: 타석 시작/중지/결과(핵심 수동)",
            "• 사용자: 금액 선택 배팅 + 고정 배당 + 광고 보상",
            "• API-SPORTS: 스코어보드·종료 상태만 자동화",
            "• 한계: 문자중계/타석 결과 없음 → 완전 자동 예측 불가",
            "• 가장 큰 운영 리스크: API 미연결·경기 매핑 오류·미정산 종료",
        ],
        top=1.7,
        size=17,
        gap=0.5,
    )
    add_footer(s, "문서 파일: docs/PPAMONG_예측게임_시스템흐름.pptx")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
