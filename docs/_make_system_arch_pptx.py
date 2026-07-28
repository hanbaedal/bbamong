# -*- coding: utf-8 -*-
"""ppamong 시스템 구조 PPT 생성 (2026-07)"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "PPAMONG_시스템_구조도.pptx"

C_TITLE = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT = RGBColor(0xE1, 0x19, 0x36)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CLIENT = RGBColor(0x1E, 0x40, 0xAF)
C_HOST = RGBColor(0x0F, 0x76, 0x6E)
C_MALL = RGBColor(0xB4, 0x53, 0x09)
C_DATA = RGBColor(0x5B, 0x21, 0xB6)
C_ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)
FONT = "Malgun Gothic"


def set_run(run, text, size=18, bold=False, color=C_TITLE):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text, 10, color=C_MUTED)


def slide_title(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, title, 28, True, C_WHITE)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run()
        set_run(sr, subtitle, 14, color=C_MUTED)
    return slide


def add_bullets(slide, items, left=0.55, top=1.55, width=12.2, size=15, gap=0.38):
    y = top
    for item in items:
        box = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(gap))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        set_run(r, item, size)
        y += gap


def add_table(slide, headers, rows, left=0.55, top=1.55, col_widths=None):
    cols = len(headers)
    rows_n = len(rows) + 1
    if col_widths is None:
        total = 12.2
        col_widths = [total / cols] * cols

    table_shape = slide.shapes.add_table(
        rows_n,
        cols,
        Inches(left),
        Inches(top),
        Inches(sum(col_widths)),
        Inches(0.45 + 0.42 * rows_n),
    ).table

    for ci, w in enumerate(col_widths):
        table_shape.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table_shape.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TITLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = C_WHITE

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table_shape.cell(ri, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_ROW_ALT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(10)
                    r.font.color.rgb = C_TITLE


def add_boxes(slide, boxes, top=1.6):
    n = len(boxes)
    gap = 0.15
    width = (12.2 - gap * (n - 1)) / n
    x = 0.55
    for label, sub, color in boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top),
            Inches(width),
            Inches(1.35),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, label, 13, True, C_WHITE)
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            set_run(r2, sub, 10, color=C_WHITE)
        x += width + gap


def add_flow_row(slide, labels, colors, y, arrow=True):
    n = len(labels)
    gap = 0.12
    w = (12.2 - gap * (n - 1)) / n
    x = 0.55
    for i, (label, color) in enumerate(zip(labels, colors)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(0.72),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, label, 11, True, C_WHITE)
        x += w + gap
    if arrow:
        arr = slide.shapes.add_textbox(Inches(5.9), Inches(y + 0.85), Inches(1.5), Inches(0.3))
        ar = arr.text_frame.paragraphs[0].add_run()
        set_run(ar, "▼", 16, True, C_MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 표지
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_TITLE
    bg.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.2))
    r = t.text_frame.paragraphs[0].add_run()
    set_run(r, "PPAMONG 시스템 구조", 40, True, C_WHITE)
    st = s.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(11.5), Inches(1.5))
    lines = [
        "게임 · 쇼핑몰 · 관리자 백오피스 통합 아키텍처",
        "2026년 7월 기준 (Replit · MongoDB Atlas · Redis)",
    ]
    for i, line in enumerate(lines):
        p = st.text_frame.paragraphs[0] if i == 0 else st.text_frame.add_paragraph()
        sr = p.add_run()
        set_run(sr, line, 16, color=RGBColor(0xCB, 0xD5, 0xE1))

    # 2 전체 개요
    s = slide_title(prs, "1. 전체 시스템 개요", "백엔드 1대(Replit) · 클라이언트 4종 · DB MongoDB Atlas")
    add_flow_row(
        s,
        ["사용자 앱", "매니저 앱", "쇼핑몰 웹", "관리자 웹"],
        [C_CLIENT, C_HOST, C_MALL, C_ACCENT],
        1.55,
    )
    add_flow_row(
        s,
        ["가비아 DNS\nppamong.com", "Replit Autoscale\nNode.js + Express", "MongoDB Atlas\n+ Redis"],
        [C_DATA, C_HOST, C_DATA],
        2.65,
        arrow=False,
    )
    add_bullets(
        s,
        [
            "• main.tsx 경로 분기: /manager → ManagerApp · /admin → AdminApp · /shop → MallApp · 그 외 → UserApp",
            "• Capacitor 앱(사용자·매니저)은 WebView로 ppamong.com 접속 — 서버 배포만으로 UI 대부분 반영",
            "• OAuth(카카오·구글·애플), SOLAPI(SMS), Object Storage(영상·상품 이미지) 연동",
        ],
        top=3.65,
        size=14,
    )
    add_footer(s, "docs/PPAMONG_시스템_구조도.md · GitHub hanbaedal/bbamong")

    # 3 클라이언트 역할
    s = slide_title(prs, "2. 클라이언트 4종 · 역할", "URL · 사용자 · 주요 기능")
    add_table(
        s,
        ["클라이언트", "접속", "대상", "주요 기능"],
        [
            ["사용자 앱", "/login, /home, /prediction …", "참여자", "게임 · 회원가입(유일) · 몰 링크"],
            ["쇼핑몰 웹", "/shop, /shop/product/:id", "일반 고객", "상품·장바구니 · 정회원 주문 · 구매 문의"],
            ["매니저 앱", "/manager/*", "현장 운영자", "경기 운영 · 예측 진행 (관리자 웹 불가)"],
            ["관리자 웹", "/admin/*", "슈퍼·일반 어드민", "게임·회원·몰·주문·재고·구매 백오피스"],
        ],
        col_widths=[2.2, 3.2, 1.8, 5.0],
    )
    add_footer(s, "정회원만 쇼핑몰 주문 가능 — 게스트·비로그인은 열람·장바구니만")

    # 4 서버 계층
    s = slide_title(prs, "3. 서버 · API 계층", "Express 단일 프로세스 · UserRoutes 모듈화")
    add_boxes(
        s,
        [
            ("게임 API", "/api/matches\n/api/predictions", C_CLIENT),
            ("쇼핑몰 API", "/api/mall/*\n상품·주문·문의", C_MALL),
            ("관리자 API", "/api/admin/*\n몰·회원·경기", C_ACCENT),
            ("공통", "WebSocket\nRedis 세션", C_HOST),
        ],
        top=1.5,
    )
    add_bullets(
        s,
        [
            "• 빌드: npm run build → Vite(React) + esbuild(server) · 실행: npm run start (PORT 5000)",
            "• 정적 파일: dist/public · 업로드: /uploads (로컬) 또는 GCS Object Storage",
            "• 배포: Replit Autoscale · git fetch && reset --hard origin/main && npm run build",
        ],
        top=3.2,
        size=13,
    )

    # 5 쇼핑몰 고객 흐름
    s = slide_title(prs, "4. 쇼핑몰(/shop) 고객 흐름", "굿웨어몰형 UI · 앱 계정과 세션 공유")
    add_flow_row(
        s,
        ["카테고리·검색", "상품 상세", "장바구니", "주문(정회원)", "구매 문의"],
        [C_MALL, C_MALL, C_MALL, C_ACCENT, C_CLIENT],
        1.55,
    )
    add_table(
        s,
        ["항목", "정책"],
        [
            ["운영 URL", "https://ppamong.com/shop (향후 shop.ppamong.com)"],
            ["회원가입", "사용자 앱 전용 — 몰 웹에 가입 폼 없음"],
            ["결제 1차", "현금 주문 접수 → 관리자 확인"],
            ["판매 유형", "재고판매(stock) · 주문후조달(procure)"],
            ["문의", "상품 상세 탭 · 이름·연락처·내용 → 관리자 구매 문의함"],
        ],
        top=2.55,
        col_widths=[2.5, 9.7],
    )

    # 6 몰 관리자 — 쇼핑몰 관리 탭
    s = slide_title(prs, "5. 쇼핑몰 관리 (/admin/mall-management)", "4개 탭 · 상품 등록/리스트/표시/문의")
    add_boxes(
        s,
        [
            ("상품 등록", "신규 POST\n등록 후 초기화", C_MALL),
            ("상품 리스트", "검색·수정\n삭제", C_MALL),
            ("쇼핑몰 표시", "제목·노출\n문의 연락처", C_HOST),
            ("구매 문의", "답변 저장\n전화·메일", C_ACCENT),
        ],
        top=1.45,
    )
    add_bullets(
        s,
        [
            "• 상품: 카테고리 · 옵션(컬러/사이즈) · 재고판매/주문후조달 · 제품사진·상품정보 이미지",
            "• 이미지: 대표 20KB/1280px · 상품정보 80KB/860px · WebP 자동압축 · GCS 또는 로컬 fallback",
            "• 구매 문의: 답변글 DB 저장 · tel:/mailto: 링크 · 처리완료 표시",
        ],
        top=3.15,
        size=13,
    )

    # 7 몰 운영 4모듈
    s = slide_title(prs, "6. 몰 운영 4모듈 (백오피스)", "주문 · 판매 · 재고 · 구매 — ERP형 흐름")
    add_table(
        s,
        ["모듈", "경로", "역할"],
        [
            ["주문 관리", "/admin/mall-orders", "고객 주문 접수·상태·배송 처리"],
            ["판매 관리", "/admin/mall-sales", "판매 실적·출고 연동"],
            ["재고 관리", "/admin/mall-inventory", "창고·로케이션·입출고·재고 수량"],
            ["구매 관리", "/admin/mall-purchase", "주문후조달 발주·입고·발송"],
        ],
        col_widths=[2.0, 3.5, 6.7],
        top=1.55,
    )
    add_bullets(
        s,
        [
            "• 재고판매: 주문 시 재고 검증·차감 · 품절 시 「판매완료」",
            "• 주문후조달: 재고 숫자 숨김 · 주문 → 구매관리 발주 → 재고 입고 → 주문 발송",
        ],
        top=4.0,
        size=14,
    )

    # 8 관리자 메뉴
    s = slide_title(prs, "7. 관리자 웹 메뉴 구조", "adminMenuConfig.ts · 슈퍼어드민 전용 구역 포함")
    add_table(
        s,
        ["섹션", "대표 메뉴"],
        [
            ["기본", "대시보드 · 앱 홈 설정"],
            ["쇼핑몰", "몰 확인 · 몰 관리 · 주문·판매·재고·구매"],
            ["슈퍼바이저", "관리자 등록/리스트 · DB백업 · 로그인 현황"],
            ["수익·운영자", "동영상 광고 수익 · 운영자 등록·모니터링"],
            ["경기·회원", "경기 관리 · 회원 리스트 · 사회공헌"],
            ["공지·지원", "공지 · 고객 지원 · 약관"],
        ],
        col_widths=[2.2, 10.0],
        top=1.5,
    )

    # 9 데이터 · 인프라
    s = slide_title(prs, "8. 데이터 · 인프라", "MongoDB · Redis · Object Storage · DNS")
    add_table(
        s,
        ["구성", "내용"],
        [
            ["MongoDB Atlas", "회원·경기·상품(GoodsProduct)·주문·ShopInquiry·리뷰 등"],
            ["Redis", "세션 · 캐시 (Replit 내장)"],
            ["Object Storage", "PRIVATE_OBJECT_DIR · 영상·광고·몰 상품 이미지 (GCS)"],
            ["로컬 fallback", "Object Storage 미설정 시 data/uploads/mall-products/"],
            ["DNS", "가비아 → ppamong.com → Replit"],
            ["타임존", "Asia/Seoul (KST)"],
        ],
        col_widths=[2.8, 9.4],
        top=1.5,
    )

    # 10 폴더 구조
    s = slide_title(prs, "9. 저장소 폴더 구조", "monorepo web — client + server + shared")
    add_bullets(
        s,
        [
            "web/",
            "  client/src/     UserApp · ManagerApp · AdminApp · MallApp · pages/mall · adminPages",
            "  server/         Express · UserRoutes · UserStorage · mongodb/models",
            "  shared/         mallConfig · mallProduct · Zod 공유 상수",
            "  android-manager/ · ios-manager-standalone/   매니저 네이티브",
            "  docs/           정책·구조·배포 문서",
        ],
        top=1.5,
        size=13,
        gap=0.32,
    )

    # 11 관련 문서
    s = slide_title(prs, "10. 관련 문서 · 배포", "운영 시 참고")
    add_table(
        s,
        ["문서", "내용"],
        [
            ["PPAMONG_시스템_구조도.md", "Mermaid 다이어그램 · 본 PPT 원본"],
            ["PPAMONG_몰_정책.md", "URL · 정회원 · 결제 · 역할"],
            ["PPAMONG_몰_판매유형.md", "stock / procure 상세"],
            ["PPAMONG_DEPLOY_CHECKLIST.md", "Replit 배포 체크리스트"],
            [".env.example", "Object Storage · SOLAPI Secrets"],
        ],
        col_widths=[4.0, 8.2],
        top=1.55,
    )
    add_bullets(
        s,
        ["Replit 배포: git fetch origin && git reset --hard origin/main && npm run build"],
        top=4.35,
        size=14,
    )

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
